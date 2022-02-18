from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
import os
from os import listdir
import PySimpleGUI as gui
from tkinter import *
from tkinter import messagebox
from tkinter import filedialog

choose_images_column = [
    [
        gui.Text("Image Folder"),
        gui.In(size=(25, 1), enable_events=True, key="-FOLDER-"),
        gui.FolderBrowse(),
    ],
]

settings_column = [
    [gui.Text("Choose your settings below")],
    [gui.Text('Choose your slide width (in inches)'), gui.InputText(key='-slidewidth-')],
    [gui.Text('Choose your slide height (in inches)'), gui.InputText(key='-slideheight-')],
    [gui.Text('Choose your left/right margins (in inches)'), gui.InputText(key='-slideleftright-')],
    [gui.Text('Choose your top/bottom margins (in inches)'), gui.InputText(key='-slidetopbottom-')],
    [gui.Button('Ok')],

]

layout = [
    [
        gui.Column(choose_images_column),
        gui.VSeperator(),
        gui.Column(settings_column),
        
    ]
]

window = gui.Window("Choose your images to import into PowerPoint", layout)


while True:
    event, values = window.read()
    if event == "Exit" or event == gui.WIN_CLOSED:
        os._exit(0)
        break
    elif event == 'Ok':
        folder_dir = values["-FOLDER-"]
        slidewidth = values["-slidewidth-"]
        slideheight = values["-slideheight-"]
        slidetopbottom = values["-slidetopbottom-"]
        slideleftright = values["-slideleftright-"]
        
        fl = values['-FOLDER-']
        sw = values['-slidewidth-']
        sh = values['-slideheight-']
        stp = values['-slidetopbottom-']
        slr = values['-slideleftright-']

        confirm = 'true'
        
        if fl == '':
            gui.Popup('You need to select a folder to import your images from!')
            confirm = 'false'
        
        if sw == '':
            gui.Popup('You need to enter a number for slide width!')
            confirm = 'false'
        else:
            try:
                value = float(sw)
            except:
                gui.Popup('You can only use numbers for slide width')
                confirm = 'false'
                
        if sh == '':
            gui.Popup('You need to enter a number for slide height!')
            confirm = 'false'
        else:
            try:
                value = float(sh)
            except:
                gui.Popup('You can only use numbers for slide height')
                confirm = 'false'
                
        if stp == '':
            gui.Popup('You need to enter a number for slide top/bottom margins!')
            confirm = 'false'
        else:
            try:
                value = float(stp)
            except:
                gui.Popup('You can only use numbers for slide top/bottom margins')
                confirm = 'false'

        if slr == '':
            gui.Popup('You need to enter a number for slide left/right margins!')
            confirm = 'false'
        else:
            try:
                value = float(slr)
            except:
                gui.Popup('You can only use numbers for slide left/right margins')
                confirm = 'false'

        if confirm == 'true':
            break


window.close()
print(sw)
print(sh)
print(stp)
print(slr)

presentation = Presentation()
presentation.slide_width = Inches(float(slidewidth))
presentation.slide_height = Inches(float(slideheight))
leftright = Inches(float(slideleftright))
topbottom = Inches(float(slidetopbottom))


for images in os.listdir(folder_dir):

    if (images.endswith(".png") or images.endswith(".jpg")\
        or images.endswith(".jpeg") or images.endswith(".PNG")):
        layout = presentation.slide_masters[0].slide_layouts[6]
        slide = presentation.slides.add_slide(layout)

        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(folder_dir + "\\" + images, left, top)

        heightIfWidthUsed = (presentation.slide_width - leftright) * pic.height / pic.width
        widthIfHeightUsed = (presentation.slide_height - topbottom) * pic.width / pic.height

        if heightIfWidthUsed > (presentation.slide_height - topbottom):

            pic.width = widthIfHeightUsed
            pic.height = (presentation.slide_height - topbottom)

        elif heightIfWidthUsed == (presentation.slide_height - topbottom):

            pic.width = (presentation.slide_width - leftright)
            pic.height = (presentation.slide_height - topbottom)

        else:
            if (presentation.slide_width - leftright) < (presentation.slide_height - topbottom):
                pic.width = (presentation.slide_width - leftright)
                pic.height = (presentation.slide_width - leftright)
            else:
                pic.width = (presentation.slide_width - leftright)
                pic.height = (presentation.slide_height - topbottom)

    
        
        #pic.width = float(presentation.slide_width - leftright)
        #pic.height = float(presentation.slide_width - topbottom)
        pic.left = int((presentation.slide_width - pic.width) / 2)
        pic.top = int((presentation.slide_height - pic.height) / 2)
        print(pic.top)

filename = filedialog.asksaveasfilename(filetypes=[("PowerPoint Presentation", "*.pptx")])
filetype = ".pptx"
saveaspptx = filename+filetype
presentation.save(saveaspptx)

App = Tk() 
App.withdraw()


messagebox.showinfo('Completed!', 'Your import of Images into PowerPoint has completed successfully!')
os._exit(0)
App.mainloop()
    
